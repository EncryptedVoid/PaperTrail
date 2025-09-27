"""
Metadata Pipeline Module

A robust file processing pipeline that handles comprehensive metadata extraction
from various file types including images, documents, PDFs, office files, audio,
video, archives, and more.

This module provides functionality to extract metadata by:
- Detecting file types and routing to appropriate extractors
- Extracting filesystem metadata (size, dates, permissions)
- Processing image EXIF, IPTC, XMP data, color profiles, and technical details
- Extracting document properties, structure information, and content analysis
- Processing audio/video metadata including codecs, duration, and technical specs
- Analyzing archive contents and compression details
- Extracting code and text file metrics and language detection
- Handling extraction failures gracefully with fallback methods
- Maintaining detailed operation logs and error tracking
- Updating artifact profiles with extracted metadata

Author: Ashiq Gazi
"""

import logging
import json
import mimetypes
import hashlib
import re
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Any, TypedDict, Optional
from tqdm import tqdm
from config import ARTIFACT_PROFILES_DIR, ARTIFACT_PREFIX, PROFILE_PREFIX


class TextCodeInfo(TypedDict):
    """Type definition for text/code analysis results"""

    code_lines: int
    comment_lines: int
    blank_lines: int
    total_lines: int
    comment_ratio: float
    functions_detected: int
    classes_detected: int
    language: str


class MetadataReport(TypedDict):
    """
    Type definition for the metadata extraction report returned by the extract_metadata method.

    Attributes:
        processed_files: Number of files that had metadata successfully extracted
        total_files: Total number of files discovered in the source directory
        failed_extractions: Count of files where metadata extraction failed
        skipped_files: Count of files skipped due to validation issues
        image_files_processed: Count of image files that were processed
        document_files_processed: Count of document files that were processed
        errors: List of error messages encountered during processing
        extraction_summary: Dictionary summarizing extraction results by file type
        profile_updates: Number of artifact profiles successfully updated
        code_analysis: Analysis results for code files
    """

    processed_files: int
    total_files: int
    failed_extractions: int
    skipped_files: int
    image_files_processed: int
    document_files_processed: int
    errors: List[str]
    extraction_summary: Dict[str, int]
    profile_updates: int
    code_analysis: Optional[TextCodeInfo]


class MetadataPipeline:
    """
    An metadata extraction pipeline for processing directories of artifacts.

    This class handles the extraction of comprehensive metadata from various file types,
    including filesystem information, image EXIF data, document properties, audio/video
    metadata, archive analysis, and much more. It maintains artifact profiles and
    provides detailed reporting of all operations.

    The pipeline works in the following stages:
    1. Directory validation and file discovery
    2. File type detection and routing
    3. Filesystem metadata extraction
    4. Specialized metadata extraction (images, documents, audio, video, etc.)
    5. Profile data integration and updates
    6. Error handling and fallback processing
    7. Comprehensive reporting and logging
    """

    def __init__(self, logger: logging.Logger):
        """
        Initialize the MetadataPipeline.

        Args:
            logger: Logger instance for recording pipeline operations and errors
        """
        self.logger = logger

    def extract_metadata(
        self, source_dir: Path, review_dir: Path, success_dir: Path
    ) -> MetadataReport:
        """
        Extract metadata from all files in a directory and update their profiles.

        This method performs comprehensive metadata extraction:
        1. Validates input directories exist and are accessible
        2. Discovers all artifact files in the source directory
        3. Loads existing artifact profiles for metadata integration
        4. Extracts filesystem and specialized metadata for each file
        5. Updates artifact profiles with extracted metadata
        6. Moves processed files to target directory
        7. Generates comprehensive report of all operations

        Args:
            source_dir: Directory containing files to extract metadata from
            success_dir: Directory to move files to after metadata extraction

        Returns:
            MetadataReport containing detailed results of the metadata extraction process

        Note:
            This method expects files to follow the ARTIFACT-{uuid}.ext naming convention
            and looks for corresponding PROFILE-{uuid}.json files for updates.
        """

        # Validate that input directories exist and are accessible
        if not source_dir.exists():
            error_msg = f"Source directory does not exist: {source_dir}"
            self.logger.error(error_msg)
            raise FileNotFoundError(error_msg)

        if not success_dir.exists():
            success_dir.mkdir(parents=True, exist_ok=True)
            self.logger.info(f"Created target directory: {success_dir}")

        # Initialize report structure with default values
        report: MetadataReport = {
            "processed_files": 0,
            "total_files": 0,
            "failed_extractions": 0,
            "skipped_files": 0,
            "image_files_processed": 0,
            "document_files_processed": 0,
            "errors": [],
            "extraction_summary": {},
            "profile_updates": 0,
            "code_analysis": None,
        }

        # Discover all artifact files in the source directory
        try:
            unprocessed_artifacts = [
                item
                for item in source_dir.iterdir()
                if item.is_file() and item.name.startswith(f"{ARTIFACT_PREFIX}-")
            ]
        except Exception as e:
            error_msg = f"Failed to scan source directory: {e}"
            self.logger.error(error_msg)
            report["errors"].append(error_msg)
            return report

        # Handle empty directory case
        if not unprocessed_artifacts:
            self.logger.info("No artifact files found in source directory")
            return report

        # Sort files by size for consistent processing order
        unprocessed_artifacts.sort(key=lambda p: p.stat().st_size)
        report["total_files"] = len(unprocessed_artifacts)

        # Log metadata extraction stage header for clear progress tracking
        self.logger.info("=" * 80)
        self.logger.info("METADATA EXTRACTION AND PROFILE UPDATE STAGE")
        self.logger.info("=" * 80)
        self.logger.info(
            f"Found {len(unprocessed_artifacts)} artifact files to process"
        )

        # Process each file through the sanitization pipeline
        # Use progress bar only for multiple files
        if len(unprocessed_artifacts) > 1:
            artifact_iterator: Any = tqdm(
                unprocessed_artifacts,
                desc="Extracting technical metadata",
                unit="artifacts",
            )
        else:
            artifact_iterator = unprocessed_artifacts

        # Process each artifact through the metadata extraction pipeline
        for artifact in artifact_iterator:
            try:
                # STAGE 1: Extract UUID from filename for profile lookup
                artifact_id = artifact.stem[
                    (len(ARTIFACT_PREFIX) + 1) :
                ]  # Remove "ARTIFACT-" prefix
                profile_path = (
                    ARTIFACT_PROFILES_DIR / f"{PROFILE_PREFIX}-{artifact_id}.json"
                )

                # STAGE 2: Load existing profile
                if not profile_path.exists():
                    error_msg = f"Profile not found for artifact: {artifact.name}"
                    self.logger.error(error_msg)
                    report["errors"].append(error_msg)
                    report["skipped_files"] += 1

                    # Move file to review directory for manual inspection
                    try:
                        review_location = review_dir / artifact.name
                        artifact.rename(review_location)
                        self.logger.info(
                            f"Moved file with missing profile to review: {review_location}"
                        )
                    except Exception as move_error:
                        self.logger.error(
                            f"Failed to move {artifact.name} to review directory: {move_error}"
                        )
                    continue

                with open(profile_path, "r", encoding="utf-8") as f:
                    profile_data = json.load(f)

                # STAGE 3: Extract metadata using the extraction engine
                self.logger.info(f"Extracting metadata for: {artifact.name}")
                extracted_metadata: Dict[str, Any] = self._extract_file_metadata(
                    artifact
                )

                # STAGE 4: Check extraction success and update counters
                if "error" in extracted_metadata:
                    report["failed_extractions"] += 1
                    report["errors"].append(
                        f"Metadata extraction failed for {artifact.name}: {extracted_metadata['error']}"
                    )
                else:
                    # Update type-specific counters
                    if extracted_metadata.get("type") == "image":
                        report["image_files_processed"] += 1
                    elif extracted_metadata.get("type") == "document":
                        report["document_files_processed"] += 1

                    # Update extraction summary by file extension
                    ext = artifact.suffix.lower()
                    report["extraction_summary"][ext] = (
                        report["extraction_summary"].get(ext, 0) + 1
                    )

                # STAGE 5: Update profile with metadata and stage completion
                profile_data["metadata"] = extracted_metadata
                profile_data["stages"] = profile_data.get("stages", {})
                profile_data["stages"]["metadata_extraction"] = {
                    "status": "completed",
                    "timestamp": datetime.now().isoformat(),
                }

                # Save updated profile
                with open(profile_path, "w", encoding="utf-8") as f:
                    json.dump(profile_data, f, indent=2)

                report["profile_updates"] += 1

                # STAGE 6: Move artifact to target directory
                moved_artifact = artifact.rename(success_dir / artifact.name)
                report["processed_files"] += 1

                self.logger.debug(f"Metadata extracted and saved for: {artifact.name}")

            except Exception as e:
                error_msg = f"Failed to process {artifact.name}: {e}"
                self.logger.error(error_msg)
                report["errors"].append(error_msg)
                report["failed_extractions"] += 1

                # Move failed file to failed subdirectory
                failed_dir = source_dir / "failed"
                failed_dir.mkdir(exist_ok=True)

                try:
                    artifact.rename(failed_dir / artifact.name)
                    self.logger.info(
                        f"Moved failed file to: {failed_dir / artifact.name}"
                    )
                except Exception as move_error:
                    self.logger.error(
                        f"Failed to move {artifact.name} to failed directory: {move_error}"
                    )

        # Generate final summary report for user review
        self.logger.info("Metadata extraction complete:")
        self.logger.info(
            f"  - {report['processed_files']} files successfully processed"
        )
        self.logger.info(f"  - {report['image_files_processed']} image files processed")
        self.logger.info(
            f"  - {report['document_files_processed']} document files processed"
        )
        self.logger.info(f"  - {report['profile_updates']} profiles updated")
        self.logger.info(f"  - {report['failed_extractions']} extraction failures")
        self.logger.info(f"  - {report['skipped_files']} files skipped")

        # Warn about any errors encountered during processing
        if report["errors"]:
            self.logger.warning(f"  - {len(report['errors'])} errors encountered")

        return report

    def _extract_file_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Entry point for metadata extraction with expanded file type support"""
        self.logger.info(f"Starting metadata extraction for: {filepath}")

        if not filepath.exists():
            self.logger.error(f"File not found: {filepath}")
            return {"error": "File not found", "path": str(filepath)}

        # Always get basic info and security analysis
        self.logger.debug("Extracting basic filesystem metadata")
        result = self._get_basic_info(filepath)

        if "error" in result:
            self.logger.error(
                f"Failed to extract basic info for {filepath}: {result['error']}"
            )
            return result

        # Route to specific extractors based on file extension
        ext = filepath.suffix.lower()
        self.logger.debug(f"File extension detected: {ext}")

        # Image files
        if ext in {
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".tiff",
            ".tif",
            ".webp",
            ".heic",
            ".heif",
            ".raw",
            ".cr2",
            ".nef",
            ".arw",
            ".dng",
            ".orf",
            ".rw2",
            ".pef",
            ".raf",
            ".3fr",
            ".dcr",
            ".kdc",
            ".srf",
            ".sr2",
            ".erf",
            ".mef",
            ".mrw",
            ".nrw",
            ".ptx",
            ".pxn",
            ".r3d",
            ".rwl",
            ".iiq",
            ".cap",
            ".dcs",
            ".fff",
            ".mdc",
            ".mos",
            ".pcd",
            ".pfm",
            ".pgm",
            ".ppm",
            ".tga",
            ".icns",
            ".ico",
            ".cur",
            ".ani",
        }:
            self.logger.info(f"Processing as image file: {ext}")
            result.update(self._extract_image_metadata(filepath))

        # Document files
        elif ext in {
            ".pdf",
            ".docx",
            ".xlsx",
            ".pptx",
            ".doc",
            ".xls",
            ".ppt",
            ".odt",
            ".ods",
            ".odp",
            ".rtf",
            ".pages",
            ".numbers",
            ".key",
        }:
            self.logger.info(f"Processing as document file: {ext}")
            result.update(self._extract_document_metadata(filepath))

        # Audio files (new support)
        elif ext in {
            ".mp3",
            ".wav",
            ".flac",
            ".m4a",
            ".aac",
            ".ogg",
            ".wma",
            ".aiff",
            ".au",
            ".ra",
            ".amr",
            ".opus",
            ".3gp",
            ".ac3",
            ".dts",
            ".ape",
            ".tak",
            ".tta",
            ".wv",
            ".mka",
            ".mp2",
            ".mp1",
            ".mpa",
            ".m1a",
            ".m2a",
            ".m4p",
            ".m4b",
            ".m4r",
            ".oga",
            ".spx",
            ".gsm",
        }:
            self.logger.info(f"Processing as audio file: {ext}")
            result.update(self._extract_audio_metadata(filepath))

        # Video files (new support)
        elif ext in {
            ".mp4",
            ".avi",
            ".mkv",
            ".mov",
            ".wmv",
            ".flv",
            ".webm",
            ".m4v",
            ".3gp",
            ".asf",
            ".f4v",
            ".f4p",
            ".f4a",
            ".f4b",
            ".vob",
            ".ogv",
            ".drc",
            ".mng",
            ".qt",
            ".yuv",
            ".rm",
            ".rmvb",
            ".viv",
            ".amv",
            ".m4p",
            ".mpg",
            ".mp2",
            ".mpeg",
            ".mpe",
            ".mpv",
            ".m2v",
            ".svi",
            ".3g2",
            ".mxf",
            ".roq",
            ".nsv",
            ".divx",
            ".xvid",
            ".ts",
            ".mts",
        }:
            self.logger.info(f"Processing as video file: {ext}")
            result.update(self._extract_video_metadata(filepath))

        # Archive files (new support)
        elif ext in {
            ".zip",
            ".rar",
            ".7z",
            ".tar",
            ".gz",
            ".bz2",
            ".xz",
            ".lzma",
            ".z",
            ".tgz",
            ".tbz2",
            ".txz",
            ".tlz",
            ".tar.gz",
            ".tar.bz2",
            ".tar.xz",
            ".tar.lzma",
            ".ace",
            ".arj",
            ".cab",
            ".lzh",
            ".lha",
            ".zoo",
            ".arc",
            ".pak",
            ".pk3",
            ".pk4",
            ".war",
            ".jar",
            ".ear",
            ".sar",
            ".apk",
            ".deb",
            ".rpm",
            ".dmg",
            ".iso",
            ".img",
            ".vdi",
            ".vmdk",
            ".vhd",
            ".qcow2",
        }:
            self.logger.info(f"Processing as archive file: {ext}")
            result.update(self._extract_archive_metadata(filepath))

        # Text and code files (new support)
        elif ext in {
            ".txt",
            ".md",
            ".csv",
            ".json",
            ".xml",
            ".html",
            ".htm",
            ".yaml",
            ".yml",
            ".toml",
            ".ini",
            ".cfg",
            ".conf",
            ".log",
            ".sql",
            ".py",
            ".js",
            ".ts",
            ".jsx",
            ".tsx",
            ".java",
            ".cpp",
            ".c",
            ".h",
            ".hpp",
            ".cs",
            ".php",
            ".rb",
            ".go",
            ".rs",
            ".swift",
            ".kt",
            ".scala",
            ".pl",
            ".ps1",
            ".sh",
            ".bash",
            ".bat",
            ".cmd",
            ".r",
            ".m",
            ".f90",
            ".f95",
            ".f03",
            ".f08",
            ".vb",
            ".vbs",
            ".asm",
            ".s",
            ".pas",
            ".pp",
            ".lua",
            ".tcl",
            ".awk",
            ".sed",
            ".vim",
            ".emacs",
            ".el",
            ".lisp",
            ".scm",
            ".clj",
            ".hs",
            ".ml",
            ".fs",
            ".elm",
            ".dart",
            ".julia",
            ".nim",
            ".zig",
            ".v",
            ".crystal",
            ".cr",
            ".d",
            ".ada",
            ".adb",
            ".ads",
            ".pro",
            ".pri",
            ".cmake",
            ".make",
            ".mk",
            ".dockerfile",
        }:
            self.logger.info(f"Processing as text/code file: {ext}")
            result.update(self._extract_text_code_metadata(filepath))

        # CAD and 3D files (new support)
        elif ext in {
            ".dwg",
            ".dxf",
            ".step",
            ".stp",
            ".iges",
            ".igs",
            ".obj",
            ".stl",
            ".ply",
            ".3ds",
            ".max",
            ".blend",
            ".ma",
            ".mb",
            ".c4d",
            ".lwo",
            ".lws",
            ".dae",
            ".fbx",
            ".x3d",
            ".wrl",
            ".3mf",
            ".amf",
            ".gltf",
            ".glb",
            ".usd",
            ".usda",
            ".usdc",
            ".abc",
            ".vdb",
            ".bgeo",
            ".geo",
        }:
            self.logger.info(f"Processing as CAD/3D file: {ext}")
            result.update(self._extract_cad_3d_metadata(filepath))

        # Font files (new support)
        elif ext in {
            ".ttf",
            ".otf",
            ".woff",
            ".woff2",
            ".eot",
            ".pfb",
            ".pfm",
            ".afm",
            ".bdf",
            ".pcf",
            ".snf",
            ".fon",
            ".fnt",
            ".dfont",
            ".suit",
            ".ttc",
            ".otc",
        }:
            self.logger.info(f"Processing as font file: {ext}")
            result.update(self._extract_font_metadata(filepath))

        # Executable and binary files (new support)
        elif ext in {
            ".exe",
            ".dll",
            ".so",
            ".dylib",
            ".app",
            ".deb",
            ".rpm",
            ".msi",
            ".pkg",
            ".bin",
            ".run",
            ".com",
            ".scr",
            ".pif",
            ".cpl",
            ".drv",
            ".sys",
            ".vxd",
            ".ax",
            ".ocx",
            ".tlb",
            ".olb",
            ".lib",
            ".a",
            ".framework",
        }:
            self.logger.info(f"Processing as executable/binary file: {ext}")
            result.update(self._extract_executable_metadata(filepath))

        else:
            self.logger.warning(
                f"No specialized extractor available for extension: {ext}"
            )
            result["warning"] = f"No specialized extractor for {ext}"
            # Still try to extract generic file analysis
            result.update(self._extract_generic_file_metadata(filepath))

        self.logger.info(f"Metadata extraction completed for: {filepath}")
        return result

    def _get_basic_info(self, filepath: Path) -> Dict[str, Any]:
        """Extract filesystem metadata with security analysis"""
        try:
            self.logger.debug(f"Reading file stats for: {filepath}")
            stat = filepath.stat()
            mime_type, encoding = mimetypes.guess_type(str(filepath))

            # Calculate file hashes for integrity checking
            file_hashes = self._calculate_file_hashes(filepath)

            # Extract more detailed filesystem information
            result = {
                "filename": filepath.name,
                "path": str(filepath.absolute()),
                "size_bytes": stat.st_size,
                "size_mb": round(stat.st_size / (1024 * 1024), 2),
                "size_human": self._format_file_size(stat.st_size),
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "accessed": datetime.fromtimestamp(stat.st_atime).isoformat(),
                "extension": filepath.suffix.lower(),
                "mime_type": mime_type,
                "encoding": encoding,
                "permissions": oct(stat.st_mode)[-3:],
                "owner_uid": stat.st_uid,
                "group_gid": stat.st_gid,
                "is_symlink": filepath.is_symlink(),
                "is_executable": bool(stat.st_mode & 0o111),
                "hashes": file_hashes,
                "extraction_timestamp": datetime.now().isoformat(),
            }

            # Add parent directory information
            result["parent_directory"] = str(filepath.parent)
            result["depth_from_root"] = len(filepath.parts) - 1

            self.logger.debug(
                f"Basic info extracted: size={result['size_human']}, "
                f"mime_type={mime_type}, permissions={result['permissions']}"
            )
            return result

        except Exception as e:
            self.logger.error(f"Failed to read basic info for {filepath}: {e}")
            return {"error": f"Failed to read basic info: {e}"}

    def _calculate_file_hashes(self, filepath: Path) -> Dict[str, str]:
        """Calculate multiple hash types for file integrity"""
        hashes = {}
        hash_algorithms = ["md5", "sha1", "sha256"]

        try:
            # Only calculate hashes for files under 100MB to avoid performance issues
            if filepath.stat().st_size > 100 * 1024 * 1024:
                return {"note": "File too large for hash calculation (>100MB)"}

            with open(filepath, "rb") as f:
                file_data = f.read()

            for algorithm in hash_algorithms:
                hasher = hashlib.new(algorithm)
                hasher.update(file_data)
                hashes[algorithm] = hasher.hexdigest()

        except Exception as e:
            self.logger.debug(f"Hash calculation failed for {filepath}: {e}")
            hashes["error"] = str(e)

        return hashes

    def _format_file_size(self, size_bytes: float) -> str:
        """Format file size in human-readable format"""
        size = float(size_bytes)  # Convert to float at the start
        for unit in ["B", "KB", "MB", "GB", "TB"]:
            if size < 1024.0:
                return f"{size:.1f} {unit}"
            size /= 1024.0
        return f"{size_bytes:.1f} PB"

    def _extract_image_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract comprehensive image metadata with advanced analysis"""
        result = {"type": "image"}
        self.logger.debug("Starting image metadata extraction")

        # PIL metadata with color analysis
        self.logger.debug("Attempting PIL metadata extraction")
        pil_data = self._get_pil_metadata(filepath)
        if "error" not in pil_data:
            self.logger.info("PIL metadata extraction successful")
            result.update(pil_data)
        else:
            self.logger.warning(
                f"PIL metadata extraction failed: {pil_data.get('error')}"
            )

        # Advanced metadata (EXIF/IPTC/XMP)
        self.logger.debug("Attempting advanced metadata extraction")
        advanced_data = self._get_advanced_image_metadata(filepath)
        if "error" not in advanced_data:
            self.logger.info("Advanced metadata extraction successful")
            result.update(advanced_data)
        else:
            self.logger.debug(
                f"Advanced metadata extraction failed: {advanced_data.get('error')}"
            )

        # Image analysis features
        analysis_data = self._get_image_analysis_metadata(filepath)
        if "error" not in analysis_data:
            result.update(analysis_data)

        return result

    def _get_pil_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract PIL-based image metadata with color and quality analysis"""
        try:
            self.logger.debug("Importing PIL dependencies")
            from PIL import Image, ImageStat
            from PIL.ExifTags import TAGS, GPSTAGS
            import pillow_heif

            pillow_heif.register_heif_opener()

            self.logger.debug(f"Opening image with PIL: {filepath}")
            with Image.open(filepath) as img:
                # Basic info
                image_info = {
                    "dimensions": img.size,
                    "width": img.size[0],
                    "height": img.size[1],
                    "aspect_ratio": round(img.size[0] / img.size[1], 3),
                    "megapixels": round((img.size[0] * img.size[1]) / 1000000, 1),
                    "mode": img.mode,
                    "format": img.format,
                    "has_transparency": img.mode in ("RGBA", "LA")
                    or "transparency" in img.info,
                    "color_channels": len(img.getbands()),
                    "color_bands": img.getbands(),
                }

                # Color statistics
                if img.mode in ("RGB", "RGBA", "L"):
                    try:
                        stat = ImageStat.Stat(img)
                        image_info["color_stats"] = {
                            "mean": stat.mean,
                            "median": stat.median,
                            "stddev": stat.stddev,
                            "extrema": stat.extrema,
                        }
                    except Exception as e:
                        self.logger.debug(f"Color statistics failed: {e}")

                # EXIF data with GPS parsing
                exif_data = {}
                gps_data = {}

                exif = img.getexif()
                if exif:
                    for k, v in exif.items():
                        tag_name = TAGS.get(k, f"tag_{k}")
                        exif_data[tag_name] = v

                        # Extract GPS data if present
                        if tag_name == "GPSInfo":
                            try:
                                gps_dict = {}
                                for gps_tag, gps_value in v.items():
                                    gps_tag_name = GPSTAGS.get(gps_tag, gps_tag)
                                    gps_dict[gps_tag_name] = gps_value
                                gps_data = gps_dict
                            except Exception as e:
                                self.logger.debug(f"GPS data parsing failed: {e}")

                    self.logger.debug(f"EXIF data extracted: {len(exif_data)} tags")
                else:
                    self.logger.debug("No EXIF data found")

                # Try to extract camera information
                camera_info = self._extract_camera_info(exif_data)

                return {
                    "image_info": image_info,
                    "exif": exif_data,
                    "gps_data": gps_data,
                    "camera_info": camera_info,
                }

        except ImportError as e:
            self.logger.error(f"PIL dependencies not available: {e}")
            return {"error": f"PIL not available: {e}"}
        except Exception as e:
            self.logger.error(f"PIL extraction failed for {filepath}: {e}")
            return {"error": f"PIL extraction failed: {e}"}

    def _extract_camera_info(self, exif_data: Dict[str, Any]) -> Dict[str, Any]:
        """Extract and parse camera information from EXIF data"""
        camera_info = {}

        try:
            # Basic camera information
            camera_fields = {
                "Make": "camera_make",
                "Model": "camera_model",
                "Software": "software",
                "DateTime": "capture_time",
                "ExposureTime": "exposure_time",
                "FNumber": "f_number",
                "ISO": "iso",
                "FocalLength": "focal_length",
                "Flash": "flash",
                "WhiteBalance": "white_balance",
                "MeteringMode": "metering_mode",
                "ExposureMode": "exposure_mode",
                "LensModel": "lens_model",
            }

            for exif_key, info_key in camera_fields.items():
                if exif_key in exif_data:
                    camera_info[info_key] = exif_data[exif_key]

            # Calculate 35mm equivalent focal length if possible
            if "FocalLength" in exif_data and "FocalLengthIn35mmFilm" in exif_data:
                camera_info["focal_length_35mm"] = exif_data["FocalLengthIn35mmFilm"]

        except Exception as e:
            self.logger.debug(f"Camera info extraction failed: {e}")

        return camera_info

    def _get_image_analysis_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Perform advanced image analysis using computer vision techniques"""
        analysis_result = {}

        try:
            # Try OpenCV for additional analysis
            import cv2
            import numpy as np

            self.logger.debug("Performing OpenCV image analysis")
            img = cv2.cv2.imread(str(filepath))
            if img is not None:
                # Basic image properties
                height, width, channels = img.shape

                # Color histogram
                hist_data = {}
                if channels == 3:  # BGR image
                    colors = ["blue", "green", "red"]
                    for i, color in enumerate(colors):
                        hist = cv2.cv2.calcHist([img], [i], None, [256], [0, 256])
                        hist_data[color] = {
                            "mean": float(np.mean(hist)),
                            "std": float(np.std(hist)),
                            "dominant_value": int(np.argmax(hist)),
                        }

                # Blur detection using Laplacian variance
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
                blur_score = cv2.Laplacian(gray, cv2.CV_64F).var()

                # Brightness analysis
                brightness = np.mean(gray)

                analysis_result["opencv_analysis"] = {
                    "color_histogram": hist_data,
                    "blur_score": float(blur_score),
                    "is_blurry": blur_score < 100,  # Threshold for blur detection
                    "brightness": float(brightness),
                    "brightness_category": self._categorize_brightness(brightness),
                }

        except ImportError:
            self.logger.debug("OpenCV not available for image analysis")
        except Exception as e:
            self.logger.debug(f"OpenCV image analysis failed: {e}")

        return analysis_result

    def _categorize_brightness(self, brightness: float) -> str:
        """Categorize image brightness"""
        if brightness < 50:
            return "very_dark"
        elif brightness < 100:
            return "dark"
        elif brightness < 150:
            return "normal"
        elif brightness < 200:
            return "bright"
        else:
            return "very_bright"

    def _get_advanced_image_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract advanced metadata using pyexiv2 if available"""
        try:
            self.logger.debug("Attempting pyexiv2 import")
            import pyexiv2

            self.logger.debug(f"Opening image with pyexiv2: {filepath}")
            with pyexiv2.Image(str(filepath)) as img:
                result = {
                    "advanced_metadata": {
                        "exif": img.read_exif(),
                        "iptc": img.read_iptc(),
                        "xmp": img.read_xmp(),
                    }
                }
                self.logger.info(
                    "Advanced metadata (EXIF/IPTC/XMP) extracted successfully"
                )
                return result

        except ImportError:
            self.logger.debug("pyexiv2 not available - using basic EXIF only")
            return {"info": "pyexiv2 not available - using basic EXIF only"}
        except Exception as e:
            self.logger.error(
                f"Advanced metadata extraction failed for {filepath}: {e}"
            )
            return {"error": f"Advanced metadata extraction failed: {e}"}

    def _extract_document_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Document extraction with content analysis"""
        result = {"type": "document"}
        ext = filepath.suffix.lower()

        # extractors for existing formats
        extractors = {
            ".pdf": self._extract_pdf_metadata,
            ".docx": self._extract_docx_metadata,
            ".xlsx": self._extract_xlsx_metadata,
            ".pptx": self._extract_pptx_metadata,
        }

        # Basic extractors for new formats
        basic_extractors = {
            ".doc": self._extract_legacy_doc_metadata,
            ".xls": self._extract_legacy_xls_metadata,
            ".ppt": self._extract_legacy_ppt_metadata,
            ".odt": self._extract_odt_metadata,
            ".ods": self._extract_ods_metadata,
            ".odp": self._extract_odp_metadata,
            ".rtf": self._extract_rtf_metadata,
        }

        if ext in extractors:
            self.logger.info(f"Using {ext} extractor")
            result.update(extractors[ext](filepath))
        elif ext in basic_extractors:
            self.logger.info(f"Using basic {ext} extractor")
            result.update(basic_extractors[ext](filepath))
        else:
            self.logger.error(f"No extractor available for document type: {ext}")
            result["error"] = f"No extractor for {ext}"

        return result

    def _extract_pdf_metadata(self, filepath: Path) -> Dict[str, Any]:
        """PDF metadata extraction with content analysis"""
        self.logger.debug("Starting PDF metadata extraction")

        # Try PyMuPDF first for comprehensive analysis
        pymupdf_result = self._try_pymupdf(filepath)
        if "error" not in pymupdf_result:
            self.logger.info("PDF metadata extracted successfully with PyMuPDF")
            return pymupdf_result

        # Fallback to basic extraction
        self.logger.info("PyMuPDF failed, falling back to basic extraction")
        result = self._try_pdfplumber(filepath)
        if "error" not in result:
            self.logger.info("PDF metadata extracted successfully with pdfplumber")
        else:
            self.logger.error("All PDF extraction methods failed")

        return result

    def _try_pymupdf(self, filepath: Path) -> Dict[str, Any]:
        """PyMuPDF extraction with content analysis"""
        try:
            import fitz  # PyMuPDF

            self.logger.debug(f"Opening PDF with PyMuPDF: {filepath}")
            doc = fitz.open(str(filepath))
            try:
                # Basic PDF info
                pdf_info = {
                    "page_count": doc.page_count,
                    "is_encrypted": doc.needs_pass,
                    "metadata": doc.metadata,
                    "toc": doc.get_toc(),
                }

                # content analysis
                content_analysis = self._analyze_pdf_content(doc)
                pdf_info.update(content_analysis)

                # Security analysis
                security_info = self._analyze_pdf_security(doc)
                pdf_info.update(security_info)

                # Structure analysis
                structure_info = self._analyze_pdf_structure(doc)
                pdf_info.update(structure_info)

                return {"pdf_info": pdf_info}

            finally:
                doc.close()

        except ImportError:
            return {"error": "PyMuPDF not available"}
        except Exception as e:
            self.logger.error(f"PyMuPDF extraction failed: {e}")
            return {"error": f"PyMuPDF extraction failed: {e}"}

    def _analyze_pdf_content(self, doc) -> Dict[str, Any]:
        """Analyze PDF content for text, images, and other elements"""
        try:
            total_text_length = 0
            total_images = 0
            total_links = 0
            fonts_used = set()

            for page_num in range(min(doc.page_count, 10)):  # Analyze first 10 pages
                page = doc[page_num]

                # Text analysis
                text = page.get_text()
                total_text_length += len(text)

                # Image analysis
                image_list = page.get_images()
                total_images += len(image_list)

                # Link analysis
                links = page.get_links()
                total_links += len(links)

                # Font analysis
                blocks = page.get_text("dict")
                for block in blocks.get("blocks", []):
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line.get("spans", []):
                                if "font" in span:
                                    fonts_used.add(span["font"])

            return {
                "content_analysis": {
                    "estimated_total_text_length": total_text_length,
                    "total_images": total_images,
                    "total_links": total_links,
                    "unique_fonts": len(fonts_used),
                    "fonts_used": list(fonts_used)[:20],  # Limit to first 20 fonts
                    "has_text": total_text_length > 0,
                    "text_to_image_ratio": total_text_length / max(total_images, 1),
                }
            }

        except Exception as e:
            self.logger.debug(f"PDF content analysis failed: {e}")
            return {"content_analysis_error": str(e)}

    def _analyze_pdf_security(self, doc) -> Dict[str, Any]:
        """Analyze PDF security features"""
        try:
            security_info = {
                "is_password_protected": doc.needs_pass,
                "permissions": (
                    {
                        "can_print": doc.permissions & fitz.PDF_PERM_PRINT,
                        "can_modify": doc.permissions & fitz.PDF_PERM_MODIFY,
                        "can_copy": doc.permissions & fitz.PDF_PERM_COPY,
                        "can_annotate": doc.permissions & fitz.PDF_PERM_ANNOTATE,
                    }
                    if hasattr(doc, "permissions")
                    else None
                ),
            }
            return {"security_info": security_info}

        except Exception as e:
            self.logger.debug(f"PDF security analysis failed: {e}")
            return {"security_analysis_error": str(e)}

    def _analyze_pdf_structure(self, doc) -> Dict[str, Any]:
        """Analyze PDF structure and organization"""
        try:
            structure_info = {
                "has_bookmarks": len(doc.get_toc()) > 0,
                "bookmark_count": len(doc.get_toc()),
                "pdf_version": doc.metadata.get("format", "Unknown"),
                "creation_date": doc.metadata.get("creationDate"),
                "modification_date": doc.metadata.get("modDate"),
                "producer": doc.metadata.get("producer"),
                "creator": doc.metadata.get("creator"),
            }
            return {"structure_info": structure_info}

        except Exception as e:
            self.logger.debug(f"PDF structure analysis failed: {e}")
            return {"structure_analysis_error": str(e)}

    def _extract_docx_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Word document metadata extraction"""
        try:
            from docx import Document

            self.logger.debug(f"Opening DOCX document: {filepath}")
            doc = Document(str(filepath))
            props = doc.core_properties

            # Basic metadata
            document_info = {
                "title": props.title,
                "author": props.author,
                "subject": props.subject,
                "keywords": props.keywords,
                "comments": props.comments,
                "category": props.category,
                "created": props.created.isoformat() if props.created else None,
                "modified": props.modified.isoformat() if props.modified else None,
                "last_modified_by": props.last_modified_by,
                "revision": props.revision,
                "version": props.version,
            }

            # Content analysis
            content_analysis = self._analyze_docx_content(doc)
            document_info.update(content_analysis)

            return {"document_info": document_info}

        except ImportError:
            return {"error": "python-docx not available"}
        except Exception as e:
            self.logger.error(f"DOCX extraction failed: {e}")
            return {"error": f"DOCX extraction failed: {e}"}

    def _analyze_docx_content(self, doc) -> Dict[str, Any]:
        """Analyze DOCX content structure and statistics"""
        try:
            # Count various elements
            paragraph_count = len(doc.paragraphs)
            table_count = len(doc.tables)

            # Analyze text content
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)

            text_content = "\n".join(full_text)
            word_count = len(text_content.split())
            char_count = len(text_content)

            # Analyze styles if available
            styles_used = set()
            for para in doc.paragraphs:
                if para.style:
                    styles_used.add(para.style.name)

            return {
                "content_stats": {
                    "paragraph_count": paragraph_count,
                    "table_count": table_count,
                    "word_count": word_count,
                    "character_count": char_count,
                    "styles_used": list(styles_used),
                    "has_tables": table_count > 0,
                    "estimated_reading_time_minutes": max(1, word_count // 200),
                }
            }

        except Exception as e:
            self.logger.debug(f"DOCX content analysis failed: {e}")
            return {"content_analysis_error": str(e)}

    def _extract_xlsx_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Excel metadata extraction"""
        try:
            import openpyxl

            self.logger.debug(f"Opening XLSX workbook: {filepath}")
            wb = openpyxl.load_workbook(str(filepath), data_only=True)
            props = wb.properties

            # Basic metadata
            spreadsheet_info = {
                "title": props.title,
                "creator": props.creator,
                "subject": props.subject,
                "description": props.description,
                "keywords": props.keywords,
                "category": props.category,
                "created": props.created.isoformat() if props.created else None,
                "modified": props.modified.isoformat() if props.modified else None,
                "last_modified_by": props.lastModifiedBy,
                "sheet_names": wb.sheetnames,
                "sheet_count": len(wb.worksheets),
            }

            # content analysis
            content_analysis = self._analyze_xlsx_content(wb)
            spreadsheet_info.update(content_analysis)

            return {"spreadsheet_info": spreadsheet_info}

        except ImportError:
            return {"error": "openpyxl not available"}
        except Exception as e:
            self.logger.error(f"XLSX extraction failed: {e}")
            return {"error": f"XLSX extraction failed: {e}"}

    def _analyze_xlsx_content(self, wb) -> Dict[str, Any]:
        """Analyze Excel content structure and statistics"""
        try:
            total_cells = 0
            total_formulas = 0
            total_charts = 0
            sheet_details = []

            for sheet in wb.worksheets:
                # Count used cells and formulas
                used_cells = 0
                formulas = 0

                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            used_cells += 1
                            if str(cell.value).startswith("="):
                                formulas += 1

                # Count charts
                charts = len(sheet._charts) if hasattr(sheet, "_charts") else 0

                sheet_info = {
                    "name": sheet.title,
                    "used_cells": used_cells,
                    "formulas": formulas,
                    "charts": charts,
                    "max_row": sheet.max_row,
                    "max_column": sheet.max_column,
                }

                sheet_details.append(sheet_info)
                total_cells += used_cells
                total_formulas += formulas
                total_charts += charts

            return {
                "content_stats": {
                    "total_used_cells": total_cells,
                    "total_formulas": total_formulas,
                    "total_charts": total_charts,
                    "sheet_details": sheet_details,
                    "has_formulas": total_formulas > 0,
                    "has_charts": total_charts > 0,
                }
            }

        except Exception as e:
            self.logger.debug(f"XLSX content analysis failed: {e}")
            return {"content_analysis_error": str(e)}

    def _extract_pptx_metadata(self, filepath: Path) -> Dict[str, Any]:
        """PowerPoint metadata extraction"""
        try:
            from pptx import Presentation

            self.logger.debug(f"Opening PPTX presentation: {filepath}")
            prs = Presentation(str(filepath))
            props = prs.core_properties

            # Basic metadata
            presentation_info = {
                "title": props.title,
                "author": props.author,
                "subject": props.subject,
                "keywords": props.keywords,
                "comments": props.comments,
                "category": props.category,
                "created": props.created.isoformat() if props.created else None,
                "modified": props.modified.isoformat() if props.modified else None,
                "last_modified_by": props.last_modified_by,
                "revision": props.revision,
                "slide_count": len(prs.slides),
            }

            # content analysis
            content_analysis = self._analyze_pptx_content(prs)
            presentation_info.update(content_analysis)

            return {"presentation_info": presentation_info}

        except ImportError:
            return {"error": "python-pptx not available"}
        except Exception as e:
            self.logger.error(f"PPTX extraction failed: {e}")
            return {"error": f"PPTX extraction failed: {e}"}

    def _analyze_pptx_content(self, prs) -> Dict[str, Any]:
        """Analyze PowerPoint content structure and statistics"""
        try:
            total_shapes = 0
            total_text_length = 0
            slide_layouts = []

            for slide in prs.slides:
                slide_shapes = len(slide.shapes)
                total_shapes += slide_shapes

                # Extract text from slide
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                slide_text_content = " ".join(slide_text)
                total_text_length += len(slide_text_content)

                # Get layout information
                if hasattr(slide, "slide_layout"):
                    layout_name = (
                        slide.slide_layout.name
                        if hasattr(slide.slide_layout, "name")
                        else "Unknown"
                    )
                    slide_layouts.append(layout_name)

            return {
                "content_stats": {
                    "total_shapes": total_shapes,
                    "total_text_length": total_text_length,
                    "slide_layouts": slide_layouts,
                    "unique_layouts": len(set(slide_layouts)),
                    "average_shapes_per_slide": (
                        round(total_shapes / len(prs.slides), 1) if prs.slides else 0
                    ),
                    "estimated_presentation_time_minutes": max(
                        1, len(prs.slides) * 2
                    ),  # 2 minutes per slide estimate
                }
            }

        except Exception as e:
            self.logger.debug(f"PPTX content analysis failed: {e}")
            return {"content_analysis_error": str(e)}

    def _extract_audio_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract audio file metadata"""
        result = {"type": "audio"}

        try:
            # Try mutagen for comprehensive audio metadata
            from mutagen import File as MutagenFile

            self.logger.debug(f"Extracting audio metadata with mutagen: {filepath}")
            audio_file = MutagenFile(str(filepath))

            if audio_file is not None:
                # Basic audio info - ensure proper types
                duration_raw = getattr(audio_file.info, "length", 0) or 0
                bitrate_raw = getattr(audio_file.info, "bitrate", 0) or 0
                sample_rate_raw = getattr(audio_file.info, "sample_rate", 0) or 0
                channels_raw = getattr(audio_file.info, "channels", 0) or 0

                audio_info = {
                    "duration_seconds": float(duration_raw),
                    "bitrate": int(bitrate_raw),
                    "sample_rate": int(sample_rate_raw),
                    "channels": int(channels_raw),
                    "codec": str(audio_file.mime[0]) if audio_file.mime else "unknown",
                }

                # Format duration nicely
                duration = audio_info["duration_seconds"]
                if duration > 0:
                    minutes, seconds = divmod(duration, 60)
                    hours, minutes = divmod(minutes, 60)
                    audio_info["duration_formatted"] = (
                        f"{int(hours):02d}:{int(minutes):02d}:{int(seconds):02d}"
                    )

                # Extract tags/metadata with proper string conversion
                tags = {}
                if audio_file.tags:
                    for key, value in audio_file.tags.items():
                        try:
                            key_str = str(key)
                            if isinstance(value, list):
                                if len(value) == 1:
                                    tags[key_str] = str(value[0])
                                elif len(value) > 1:
                                    # Join multiple values with comma
                                    tags[key_str] = ", ".join(str(v) for v in value)
                            else:
                                tags[key_str] = str(value)
                        except (TypeError, ValueError) as e:
                            # Skip problematic values
                            self.logger.debug(f"Skipping tag {key}: {e}")
                            continue

                result["audio_info"] = audio_info
                result["tags"] = tags

                self.logger.info(
                    f"Audio metadata extracted: {audio_info.get('duration_formatted', 'unknown duration')}"
                )

        except ImportError:
            self.logger.debug("Mutagen not available for audio metadata extraction")
            result["error"] = "Mutagen library not available"
        except Exception as e:
            self.logger.error(f"Audio metadata extraction failed: {e}")
            result["error"] = f"Audio metadata extraction failed: {e}"

        return result

    def _extract_video_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract video file metadata"""
        result = {"type": "video"}

        # Try ffprobe first (most comprehensive)
        ffprobe_result = self._try_ffprobe_video(filepath)
        if "error" not in ffprobe_result:
            result.update(ffprobe_result)
            return result

        # Fallback to OpenCV
        opencv_result = self._try_opencv_video(filepath)
        if "error" not in opencv_result:
            result.update(opencv_result)
        else:
            result.update(opencv_result)  # Include error

        return result

    def _try_ffprobe_video(self, filepath: Path) -> Dict[str, Any]:
        """Try to extract video metadata using ffprobe"""
        try:
            import subprocess
            import json

            cmd = [
                "ffprobe",
                "-v",
                "quiet",
                "-print_format",
                "json",
                "-show_format",
                "-show_streams",
                str(filepath),
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)

            if result.returncode == 0:
                data = json.loads(result.stdout)

                # Extract video stream info
                video_streams = [
                    s for s in data.get("streams", []) if s.get("codec_type") == "video"
                ]
                audio_streams = [
                    s for s in data.get("streams", []) if s.get("codec_type") == "audio"
                ]

                video_info = {}
                if video_streams:
                    v = video_streams[0]  # Use first video stream
                    video_info = {
                        "width": v.get("width"),
                        "height": v.get("height"),
                        "codec": v.get("codec_name"),
                        "fps": eval(
                            v.get("r_frame_rate", "0/1")
                        ),  # Convert fraction to float
                        "duration_seconds": float(v.get("duration", 0)),
                        "bitrate": (
                            int(v.get("bit_rate", 0)) if v.get("bit_rate") else None
                        ),
                        "pixel_format": v.get("pix_fmt"),
                    }

                    # Format duration
                    duration = video_info["duration_seconds"]
                    if duration > 0:
                        minutes, seconds = divmod(duration, 60)
                        hours, minutes = divmod(minutes, 60)
                        video_info["duration_formatted"] = (
                            f"{int(hours):02d}:{int(minutes):02d}:{int(seconds):02d}"
                        )

                audio_info = {}
                if audio_streams:
                    a = audio_streams[0]  # Use first audio stream
                    audio_info = {
                        "codec": a.get("codec_name"),
                        "sample_rate": (
                            int(a.get("sample_rate", 0))
                            if a.get("sample_rate")
                            else None
                        ),
                        "channels": a.get("channels"),
                        "bitrate": (
                            int(a.get("bit_rate", 0)) if a.get("bit_rate") else None
                        ),
                    }

                format_info = data.get("format", {})

                return {
                    "video_info": video_info,
                    "audio_info": audio_info,
                    "format_info": {
                        "format_name": format_info.get("format_name"),
                        "size_bytes": (
                            int(format_info.get("size", 0))
                            if format_info.get("size")
                            else None
                        ),
                        "duration_seconds": (
                            float(format_info.get("duration", 0))
                            if format_info.get("duration")
                            else None
                        ),
                    },
                }
            else:
                return {"error": f"ffprobe failed: {result.stderr}"}

        except (ImportError, FileNotFoundError):
            return {"error": "ffprobe not available"}
        except Exception as e:
            return {"error": f"ffprobe extraction failed: {e}"}

    def _try_opencv_video(self, filepath: Path) -> Dict[str, Any]:
        """Fallback video metadata extraction using OpenCV"""
        try:
            from cv2 import VideoCapture
            from cv2 import (
                CAP_PROP_FPS,
                CAP_PROP_FRAME_COUNT,
                CAP_PROP_FRAME_WIDTH,
                CAP_PROP_FRAME_HEIGHT,
            )

            cap = VideoCapture(str(filepath))

            fps = cap.get(CAP_PROP_FPS)
            frame_count = cap.get(CAP_PROP_FRAME_COUNT)
            width = int(cap.get(CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(CAP_PROP_FRAME_HEIGHT))
            width = int(cap.get(CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(CAP_PROP_FRAME_HEIGHT))

            duration = frame_count / fps if fps > 0 else 0

            video_info = {
                "width": width,
                "height": height,
                "fps": fps,
                "frame_count": int(frame_count),
                "duration_seconds": duration,
            }

            if duration > 0:
                minutes, seconds = divmod(duration, 60)
                hours, minutes = divmod(minutes, 60)
                video_info["duration_formatted"] = (
                    f"{int(hours):02d}:{int(minutes):02d}:{int(seconds):02d}"
                )

            cap.release()

            return {"video_info": video_info}

        except ImportError:
            return {"error": "OpenCV not available"}
        except Exception as e:
            return {"error": f"OpenCV video extraction failed: {e}"}

    def _extract_archive_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract archive file metadata"""
        result = {"type": "archive"}
        ext = filepath.suffix.lower()

        # Route to specific archive handlers
        if ext in {".zip", ".jar", ".war", ".ear", ".apk"}:
            result.update(self._extract_zip_metadata(filepath))
        elif ext in {".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tbz2", ".tar.xz", ".txz"}:
            result.update(self._extract_tar_metadata(filepath))
        elif ext == ".rar":
            result.update(self._extract_rar_metadata(filepath))
        elif ext == ".7z":
            result.update(self._extract_7z_metadata(filepath))
        else:
            result.update(self._extract_generic_archive_metadata(filepath))

        return result

    def _extract_zip_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract ZIP archive metadata"""
        try:
            import zipfile

            with zipfile.ZipFile(filepath, "r") as zf:
                file_list = zf.infolist()

                archive_info = {
                    "file_count": len(file_list),
                    "uncompressed_size": sum(f.file_size for f in file_list),
                    "compressed_size": sum(f.compress_size for f in file_list),
                    "compression_ratio": 0,
                    "files": [],
                }

                # Calculate compression ratio
                if archive_info["uncompressed_size"] > 0:
                    archive_info["compression_ratio"] = round(
                        1
                        - (
                            archive_info["compressed_size"]
                            / archive_info["uncompressed_size"]
                        ),
                        3,
                    )

                # Analyze file types
                file_extensions = {}
                for file_info in file_list[:100]:  # Limit to first 100 files
                    ext = Path(file_info.filename).suffix.lower()
                    file_extensions[ext] = file_extensions.get(ext, 0) + 1

                    archive_info["files"].append(
                        {
                            "name": file_info.filename,
                            "size": file_info.file_size,
                            "compressed_size": file_info.compress_size,
                            "modified": (
                                datetime(*file_info.date_time).isoformat()
                                if file_info.date_time
                                else None
                            ),
                        }
                    )

                archive_info["file_extensions"] = file_extensions
                archive_info["uncompressed_size_mb"] = round(
                    archive_info["uncompressed_size"] / (1024 * 1024), 2
                )

                return {"archive_info": archive_info}

        except Exception as e:
            return {"error": f"ZIP extraction failed: {e}"}

    def _extract_tar_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract TAR archive metadata"""
        try:
            import tarfile

            with tarfile.open(filepath, "r:*") as tf:
                members = tf.getmembers()

                archive_info = {
                    "file_count": len(members),
                    "total_size": sum(m.size for m in members if m.isfile()),
                    "directories": sum(1 for m in members if m.isdir()),
                    "regular_files": sum(1 for m in members if m.isfile()),
                    "symbolic_links": sum(1 for m in members if m.islnk() or m.issym()),
                    "files": [],
                }

                # Analyze file types
                file_extensions = {}
                for member in members[:100]:  # Limit to first 100 files
                    if member.isfile():
                        ext = Path(member.name).suffix.lower()
                        file_extensions[ext] = file_extensions.get(ext, 0) + 1

                        archive_info["files"].append(
                            {
                                "name": member.name,
                                "size": member.size,
                                "modified": (
                                    datetime.fromtimestamp(member.mtime).isoformat()
                                    if member.mtime
                                    else None
                                ),
                                "type": (
                                    "file"
                                    if member.isfile()
                                    else "directory" if member.isdir() else "link"
                                ),
                            }
                        )

                archive_info["file_extensions"] = file_extensions
                archive_info["total_size_mb"] = round(
                    archive_info["total_size"] / (1024 * 1024), 2
                )

                return {"archive_info": archive_info}

        except Exception as e:
            return {"error": f"TAR extraction failed: {e}"}

    def _extract_rar_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract RAR archive metadata"""
        try:
            import rarfile

            with rarfile.RarFile(filepath) as rf:
                file_list = rf.infolist()

                archive_info = {
                    "file_count": len(file_list),
                    "uncompressed_size": sum(f.file_size for f in file_list),
                    "compressed_size": sum(f.compress_size for f in file_list),
                    "files": [],
                }

                # Calculate compression ratio
                if archive_info["uncompressed_size"] > 0:
                    archive_info["compression_ratio"] = round(
                        1
                        - (
                            archive_info["compressed_size"]
                            / archive_info["uncompressed_size"]
                        ),
                        3,
                    )

                for file_info in file_list[:100]:  # Limit to first 100 files
                    archive_info["files"].append(
                        {
                            "name": file_info.filename,
                            "size": file_info.file_size,
                            "compressed_size": file_info.compress_size,
                            "modified": (
                                file_info.date_time.isoformat()
                                if hasattr(file_info, "date_time")
                                and file_info.date_time
                                else None
                            ),
                        }
                    )

                return {"archive_info": archive_info}

        except ImportError:
            return {"error": "rarfile library not available"}
        except Exception as e:
            return {"error": f"RAR extraction failed: {e}"}

    def _extract_7z_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract 7z archive metadata"""
        try:
            import py7zr

            with py7zr.SevenZipFile(filepath, mode="r") as archive:
                file_list = archive.list()

                archive_info = {"file_count": len(file_list), "files": []}

                for file_info in file_list:
                    archive_info["files"].append(
                        {
                            "name": file_info.filename,
                            "size": (
                                file_info.uncompressed
                                if hasattr(file_info, "uncompressed")
                                else None
                            ),
                            "compressed_size": (
                                file_info.compressed
                                if hasattr(file_info, "compressed")
                                else None
                            ),
                            "is_directory": (
                                file_info.is_directory
                                if hasattr(file_info, "is_directory")
                                else False
                            ),
                        }
                    )

                return {"archive_info": archive_info}

        except ImportError:
            return {"error": "py7zr library not available"}
        except Exception as e:
            return {"error": f"7z extraction failed: {e}"}

    def _extract_generic_archive_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Generic archive metadata extraction"""
        return {
            "archive_info": {
                "note": f"Generic archive handler for {filepath.suffix}",
                "file_count": "unknown",
                "extraction_method": "none_available",
            }
        }

    def _extract_text_code_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract metadata from text and code files"""
        result = {"type": "text_code"}

        try:
            # Read file content with encoding detection
            content, encoding = self._read_text_file_with_encoding(filepath)

            if content is not None:
                # Basic text statistics
                text_stats = self._analyze_text_content(
                    content, filepath.suffix.lower()
                )
                result = {
                    "type": "text_code",
                    "text_info": text_stats,
                    "encoding": encoding,
                }

                # Language detection for code files
                if filepath.suffix.lower() in {
                    ".py",
                    ".js",
                    ".java",
                    ".cpp",
                    ".c",
                    ".cs",
                    ".php",
                    ".rb",
                    ".go",
                }:
                    code_analysis = self._analyze_code_content(
                        content, filepath.suffix.lower()
                    )
                    result["code_analysis"] = code_analysis

        except Exception as e:
            result["error"] = f"Text/code analysis failed: {e}"

        return result

    def _read_text_file_with_encoding(
        self, filepath: Path
    ) -> tuple[Optional[str], str]:
        """Read text file with automatic encoding detection"""
        encodings_to_try = ["utf-8", "utf-16", "latin-1", "cp1252", "iso-8859-1"]

        for encoding in encodings_to_try:
            try:
                with open(filepath, "r", encoding=encoding) as f:
                    content = f.read()
                return content, encoding
            except UnicodeDecodeError:
                continue
            except Exception:
                break

        # Try with chardet if available
        try:
            import chardet

            with open(filepath, "rb") as f:
                raw_data = f.read()
            detected = chardet.detect(raw_data)
            if detected["encoding"]:
                with open(filepath, "r", encoding=detected["encoding"]) as f:
                    content = f.read()
                return content, detected["encoding"]
        except ImportError:
            pass
        except Exception:
            pass

        return None, "unknown"

    def _analyze_text_content(self, content: str, file_ext: str) -> Dict[str, Any]:
        """Analyze text content for various metrics"""
        lines = content.split("\n")
        words = content.split()

        # Basic statistics
        stats = {
            "line_count": len(lines),
            "word_count": len(words),
            "character_count": len(content),
            "character_count_no_spaces": len(content.replace(" ", "")),
            "empty_lines": sum(1 for line in lines if not line.strip()),
            "max_line_length": max(len(line) for line in lines) if lines else 0,
            "avg_line_length": (
                sum(len(line) for line in lines) / len(lines) if lines else 0
            ),
        }

        # Language-specific analysis
        if file_ext in {".csv"}:
            stats.update(self._analyze_csv_structure(content))
        elif file_ext in {".json"}:
            stats.update(self._analyze_json_structure(content))
        elif file_ext in {".xml", ".html", ".htm"}:
            stats.update(self._analyze_xml_html_structure(content))
        elif file_ext in {".md"}:
            stats.update(self._analyze_markdown_structure(content))

        return stats

    def _analyze_csv_structure(self, content: str) -> Dict[str, Any]:
        """Analyze CSV file structure"""
        try:
            import csv
            from io import StringIO

            # Detect dialect
            sniffer = csv.Sniffer()
            sample = content[:1024]
            dialect = sniffer.sniff(sample)

            # Parse CSV
            reader = csv.reader(StringIO(content), dialect)
            rows = list(reader)

            return {
                "csv_analysis": {
                    "delimiter": dialect.delimiter,
                    "quote_char": dialect.quotechar,
                    "row_count": len(rows),
                    "column_count": len(rows[0]) if rows else 0,
                    "has_header": sniffer.has_header(sample),
                    "first_row": rows[0] if rows else None,
                }
            }
        except Exception as e:
            return {"csv_analysis_error": str(e)}

    def _analyze_json_structure(self, content: str) -> Dict[str, Any]:
        """Analyze JSON file structure"""
        try:
            import json

            data = json.loads(content)

            def count_elements(obj, depth=0):
                if isinstance(obj, dict):
                    return {
                        "objects": 1
                        + sum(
                            count_elements(v, depth + 1).get("objects", 0)
                            for v in obj.values()
                        ),
                        "arrays": sum(
                            count_elements(v, depth + 1).get("arrays", 0)
                            for v in obj.values()
                        ),
                        "keys": len(obj),
                        "max_depth": max(
                            [depth]
                            + [
                                count_elements(v, depth + 1).get("max_depth", depth)
                                for v in obj.values()
                            ]
                        ),
                    }
                elif isinstance(obj, list):
                    return {
                        "objects": sum(
                            count_elements(item, depth + 1).get("objects", 0)
                            for item in obj
                        ),
                        "arrays": 1
                        + sum(
                            count_elements(item, depth + 1).get("arrays", 0)
                            for item in obj
                        ),
                        "items": len(obj),
                        "max_depth": max(
                            [depth]
                            + [
                                count_elements(item, depth + 1).get("max_depth", depth)
                                for item in obj
                            ]
                        ),
                    }
                else:
                    return {"max_depth": depth}

            count_analysis = count_elements(data)
            analysis = {**count_analysis, "root_type": type(data).__name__}

            return {"json_analysis": analysis}

        except Exception as e:
            return {"json_analysis_error": str(e)}

    def _analyze_xml_html_structure(self, content: str) -> Dict[str, Any]:
        """Analyze XML/HTML file structure"""
        try:
            # Simple tag counting using regex
            import re

            # Count tags
            tag_pattern = r"<(\w+)"
            tags = re.findall(tag_pattern, content, re.IGNORECASE)
            tag_counts = {}
            for tag in tags:
                tag_counts[tag.lower()] = tag_counts.get(tag.lower(), 0) + 1

            # Count attributes
            attr_pattern = r"<\w+[^>]*\s+(\w+)="
            attributes = re.findall(attr_pattern, content, re.IGNORECASE)

            return {
                "markup_analysis": {
                    "total_tags": len(tags),
                    "unique_tags": len(tag_counts),
                    "most_common_tags": sorted(
                        tag_counts.items(), key=lambda x: x[1], reverse=True
                    )[:10],
                    "total_attributes": len(attributes),
                    "has_doctype": "<!DOCTYPE" in content.upper(),
                }
            }

        except Exception as e:
            return {"markup_analysis_error": str(e)}

    def _analyze_markdown_structure(self, content: str) -> Dict[str, Any]:
        """Analyze Markdown file structure"""
        try:
            lines = content.split("\n")

            headers = []
            lists = 0
            links = 0
            images = 0
            code_blocks = 0

            for line in lines:
                line = line.strip()
                if line.startswith("#"):
                    level = len(line) - len(line.lstrip("#"))
                    headers.append(level)
                elif line.startswith(("- ", "* ", "+ ")) or (
                    line and line[0].isdigit() and ". " in line
                ):
                    lists += 1

                # Count links and images
                import re

                links += len(re.findall(r"\[.*?\]\(.*?\)", line))
                images += len(re.findall(r"!\[.*?\]\(.*?\)", line))

                if line.startswith("```"):
                    code_blocks += 1

            return {
                "markdown_analysis": {
                    "header_count": len(headers),
                    "header_levels": headers,
                    "list_items": lists,
                    "links": links,
                    "images": images,
                    "code_blocks": code_blocks
                    // 2,  # Divide by 2 since we count opening and closing
                }
            }

        except Exception as e:
            return {"markdown_analysis_error": str(e)}

    def _analyze_code_content(self, content: str, file_ext: str) -> Dict[str, Any]:
        """Analyze code content for programming metrics"""
        try:
            lines = content.split("\n")

            # Basic code metrics
            code_lines = 0
            comment_lines = 0
            blank_lines = 0

            # Comment patterns for different languages
            comment_patterns = {
                ".py": [r"#.*$"],
                ".js": [r"//.*$", r"/\*.*?\*/"],
                ".java": [r"//.*$", r"/\*.*?\*/"],
                ".cpp": [r"//.*$", r"/\*.*?\*/"],
                ".c": [r"//.*$", r"/\*.*?\*/"],
                ".cs": [r"//.*$", r"/\*.*?\*/"],
                ".php": [r"//.*$", r"/\*.*?\*/", r"#.*$"],
                ".rb": [r"#.*$"],
                ".go": [r"//.*$", r"/\*.*?\*/"],
            }

            patterns = comment_patterns.get(file_ext, [r"#.*$"])

            for line in lines:
                stripped = line.strip()
                if not stripped:
                    blank_lines += 1
                elif any(re.search(pattern, stripped) for pattern in patterns):
                    comment_lines += 1
                else:
                    code_lines += 1

            # Function/class detection (basic)
            functions = (
                len(re.findall(r"def\s+\w+\s*\(", content)) if file_ext == ".py" else 0
            )
            classes = (
                len(re.findall(r"class\s+\w+", content)) if file_ext == ".py" else 0
            )

            return {
                "code_lines": code_lines,
                "comment_lines": comment_lines,
                "blank_lines": blank_lines,
                "total_lines": len(lines),
                "comment_ratio": round(comment_lines / len(lines), 3) if lines else 0,
                "functions_detected": functions,
                "classes_detected": classes,
                "language": self._detect_programming_language(file_ext),
            }

        except Exception as e:
            return {"code_analysis_error": str(e)}

    def _detect_programming_language(self, file_ext: str) -> str:
        """Map file extension to programming language"""
        language_map = {
            ".py": "Python",
            ".js": "JavaScript",
            ".ts": "TypeScript",
            ".java": "Java",
            ".cpp": "C++",
            ".c": "C",
            ".cs": "C#",
            ".php": "PHP",
            ".rb": "Ruby",
            ".go": "Go",
            ".rs": "Rust",
            ".swift": "Swift",
            ".kt": "Kotlin",
            ".scala": "Scala",
            ".pl": "Perl",
            ".ps1": "PowerShell",
            ".sh": "Shell",
            ".bash": "Bash",
            ".r": "R",
            ".m": "MATLAB",
            ".f90": "Fortran",
            ".vb": "VB.NET",
            ".lua": "Lua",
            ".tcl": "Tcl",
            ".hs": "Haskell",
            ".ml": "OCaml",
            ".elm": "Elm",
            ".dart": "Dart",
            ".julia": "Julia",
        }
        return language_map.get(file_ext, "Unknown")

    def _extract_cad_3d_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract CAD and 3D file metadata"""
        result: Dict[str, Any] = {"type": "cad_3d"}
        ext = filepath.suffix.lower()

        if ext in {".obj"}:
            result.update(self._extract_obj_metadata(filepath))
        elif ext in {".stl"}:
            result.update(self._extract_stl_metadata(filepath))
        elif ext in {".ply"}:
            result.update(self._extract_ply_metadata(filepath))
        else:
            result["cad_3d_info"] = {
                "format": ext[1:].upper(),
                "note": f"Basic file info only for {ext} format",
            }

        return result

    def _extract_obj_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract OBJ file metadata"""
        try:
            vertices = 0
            faces = 0
            normals = 0
            textures = 0
            materials = set()

            with open(filepath, "r") as f:
                for line in f:
                    line = line.strip()
                    if line.startswith("v "):
                        vertices += 1
                    elif line.startswith("f "):
                        faces += 1
                    elif line.startswith("vn "):
                        normals += 1
                    elif line.startswith("vt "):
                        textures += 1
                    elif line.startswith("usemtl "):
                        materials.add(line.split()[1])

            return {
                "obj_info": {
                    "vertices": vertices,
                    "faces": faces,
                    "normals": normals,
                    "texture_coordinates": textures,
                    "materials": list(materials),
                    "material_count": len(materials),
                }
            }

        except Exception as e:
            return {"error": f"OBJ analysis failed: {e}"}

    def _extract_stl_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract STL file metadata"""
        try:
            # Read first 80 bytes to check if binary or ASCII
            with open(filepath, "rb") as f:
                header = f.read(80)

            is_ascii = b"solid" in header[:5].lower()

            if is_ascii:
                # ASCII STL
                triangles = 0
                with open(filepath, "r") as f:
                    for line in f:
                        if line.strip().startswith("facet normal"):
                            triangles += 1
            else:
                # Binary STL
                with open(filepath, "rb") as f:
                    f.seek(80)  # Skip header
                    triangle_count_bytes = f.read(4)
                    triangles = int.from_bytes(triangle_count_bytes, byteorder="little")

            return {
                "stl_info": {
                    "format": "ASCII" if is_ascii else "Binary",
                    "triangles": triangles,
                    "vertices": triangles * 3,  # Each triangle has 3 vertices
                }
            }

        except Exception as e:
            return {"error": f"STL analysis failed: {e}"}

    def _extract_ply_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract PLY file metadata"""
        try:
            vertices = 0
            faces = 0
            format_type = "unknown"

            with open(filepath, "r") as f:
                header_end = False
                while not header_end:
                    line = f.readline().strip()
                    if line.startswith("format"):
                        format_type = line.split()[1]
                    elif line.startswith("element vertex"):
                        vertices = int(line.split()[2])
                    elif line.startswith("element face"):
                        faces = int(line.split()[2])
                    elif line == "end_header":
                        header_end = True

            return {
                "ply_info": {
                    "format": format_type,
                    "vertices": vertices,
                    "faces": faces,
                }
            }

        except Exception as e:
            return {"error": f"PLY analysis failed: {e}"}

    def _extract_font_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract font file metadata"""
        result = {"type": "font"}

        try:
            # Try fontTools for comprehensive font analysis
            from fontTools.ttLib import TTFont

            font = TTFont(filepath)

            # Basic font info
            font_info = {
                "format": filepath.suffix[1:].upper(),
                "tables": list(font.keys()),
                "glyph_count": (
                    font.getGlyphSet().__len__()
                    if hasattr(font.getGlyphSet(), "__len__")
                    else "unknown"
                ),
            }

            # Extract name table information
            if "name" in font:
                name_table = font["name"]
                names = {}
                for record in name_table.names:
                    if record.nameID in {
                        1,
                        2,
                        4,
                        6,
                    }:  # Font family, subfamily, full name, PostScript name
                        try:
                            names[record.nameID] = record.toUnicodeString()
                        except:
                            pass
                font_info["names"] = names

            # Extract head table information (if available)
            if "head" in font:
                head = font["head"]
                font_info["created"] = (
                    head.created if hasattr(head, "created") else None
                )
                font_info["modified"] = (
                    head.modified if hasattr(head, "modified") else None
                )

            font.close()

            return {"font_info": font_info}

        except ImportError:
            return {"error": "fontTools library not available"}
        except Exception as e:
            return {"error": f"Font analysis failed: {e}"}

    def _extract_executable_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract executable and binary file metadata"""
        result = {"type": "executable"}
        ext = filepath.suffix.lower()

        try:
            # Basic executable info
            exec_info = {
                "format": ext[1:].upper(),
                "is_executable": True,
                "platform": self._detect_executable_platform(filepath, ext),
            }

            # Try to extract more detailed info for PE files (Windows)
            if ext in {".exe", ".dll"}:
                pe_info = self._extract_pe_metadata(filepath)
                if pe_info:
                    exec_info.update(pe_info)

            return {"executable_info": exec_info}

        except Exception as e:
            return {"error": f"Executable analysis failed: {e}"}

    def _detect_executable_platform(self, ext: str) -> str:
        """Detect executable platform"""
        platform_map = {
            ".exe": "Windows",
            ".dll": "Windows",
            ".msi": "Windows",
            ".so": "Linux",
            ".dylib": "macOS",
            ".app": "macOS",
            ".deb": "Debian/Ubuntu",
            ".rpm": "Red Hat/CentOS/SUSE",
        }

        return platform_map.get(ext, "Unknown")

    def _extract_pe_metadata(self, filepath: Path) -> Optional[Dict[str, Any]]:
        """Extract PE (Windows executable) metadata"""
        try:
            import pefile

            pe = pefile.PE(str(filepath))

            pe_info = {
                "machine_type": hex(pe.FILE_HEADER.Machine),
                "timestamp": datetime.fromtimestamp(
                    pe.FILE_HEADER.TimeDateStamp
                ).isoformat(),
                "sections": len(pe.sections),
                "is_dll": pe.is_dll(),
                "is_driver": pe.is_driver(),
                "is_exe": pe.is_exe(),
            }

            # Extract version info if available
            if hasattr(pe, "VS_VERSIONINFO"):
                for fileinfo in pe.VS_VERSIONINFO:
                    if hasattr(fileinfo, "StringTable"):
                        for st in fileinfo.StringTable:
                            for entry in st.entries.items():
                                pe_info[f"version_{entry[0]}"] = entry[1]

            pe.close()
            return pe_info

        except ImportError:
            return None
        except Exception:
            return None

    def _extract_generic_file_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract generic metadata for unknown file types"""
        result = {"type": "generic"}

        try:
            # Try to determine if file is binary or text
            with open(filepath, "rb") as f:
                chunk = f.read(1024)

            # Simple heuristic: if more than 30% of bytes are non-printable, consider it binary
            non_printable = sum(1 for byte in chunk if byte < 32 or byte > 126)
            is_binary = (non_printable / len(chunk)) > 0.3 if chunk else True

            generic_info = {
                "format": (
                    filepath.suffix[1:].upper() if filepath.suffix else "No extension"
                ),
                "appears_binary": is_binary,
                "sample_bytes": chunk[:50].hex() if is_binary else None,
                "sample_text": (
                    chunk[:200].decode("utf-8", errors="ignore")
                    if not is_binary
                    else None
                ),
            }

            return {"generic_info": generic_info}

        except Exception as e:
            return {"error": f"Generic analysis failed: {e}"}

    # Legacy document format extractors
    def _extract_legacy_doc_metadata(self) -> Dict[str, Any]:
        """Extract legacy DOC file metadata"""
        return {
            "document_info": {
                "format": "Legacy DOC",
                "note": "Limited metadata extraction available",
            }
        }

    def _extract_legacy_xls_metadata(self) -> Dict[str, Any]:
        """Extract legacy XLS file metadata"""
        return {
            "spreadsheet_info": {
                "format": "Legacy XLS",
                "note": "Limited metadata extraction available",
            }
        }

    def _extract_legacy_ppt_metadata(self) -> Dict[str, Any]:
        """Extract legacy PPT file metadata"""
        return {
            "presentation_info": {
                "format": "Legacy PPT",
                "note": "Limited metadata extraction available",
            }
        }

    def _extract_odt_metadata(self) -> Dict[str, Any]:
        """Extract OpenDocument Text metadata"""
        return {
            "document_info": {"format": "OpenDocument Text", "note": "Basic support"}
        }

    def _extract_ods_metadata(self) -> Dict[str, Any]:
        """Extract OpenDocument Spreadsheet metadata"""
        return {
            "spreadsheet_info": {
                "format": "OpenDocument Spreadsheet",
                "note": "Basic support",
            }
        }

    def _extract_odp_metadata(self) -> Dict[str, Any]:
        """Extract OpenDocument Presentation metadata"""
        return {
            "presentation_info": {
                "format": "OpenDocument Presentation",
                "note": "Basic support",
            }
        }

    def _extract_rtf_metadata(self) -> Dict[str, Any]:
        """Extract RTF file metadata"""
        return {
            "document_info": {"format": "Rich Text Format", "note": "Basic support"}
        }
