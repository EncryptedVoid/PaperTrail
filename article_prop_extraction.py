from pathlib import Path
from datetime import datetime
from typing import Dict, Any, Union
import mimetypes
import logging


class MetadataExtractor:
    """Comprehensive file metadata extraction with lazy loading of dependencies"""

    def __init__(self, logger: logging.Logger):
        if logger is None:
            raise ValueError(
                "Logger is required - MetadataExtractor cannot be initialized without a logger"
            )

        if not isinstance(logger, logging.Logger):
            raise TypeError(f"Expected logging.Logger instance, got {type(logger)}")

        self.logger = logger
        self._cache = {}
        self.logger.info("MetadataExtractor initialized successfully")

    def extract(self, filepath: Union[str, Path]) -> Dict[str, Any]:
        """Main entry point for metadata extraction"""
        filepath = Path(filepath)
        self.logger.info(f"Starting metadata extraction for: {filepath}")

        if not filepath.exists():
            self.logger.error(f"File not found: {filepath}")
            return {"error": "File not found", "path": str(filepath)}

        # Always get basic info
        self.logger.debug("Extracting basic filesystem metadata")
        result = self._get_basic_info(filepath)

        if "error" in result:
            self.logger.error(
                f"Failed to extract basic info for {filepath}: {result['error']}"
            )
            return result

        # Route to specific extractors
        ext = filepath.suffix.lower()
        self.logger.debug(f"File extension detected: {ext}")

        if ext in {
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".tiff",
            ".tif",
            ".webp",
            ".heic",
            ".heif",
            ".raw",
            ".cr2",
            ".nef",
            ".arw",
        }:
            self.logger.info(f"Processing as image file: {ext}")
            result.update(self._extract_image_metadata(filepath))
        elif ext in {".pdf", ".docx", ".xlsx", ".pptx"}:
            self.logger.info(f"Processing as document file: {ext}")
            result.update(self._extract_document_metadata(filepath))
        else:
            self.logger.warning(
                f"No specialized extractor available for extension: {ext}"
            )
            result["warning"] = f"No specialized extractor for {ext}"

        self.logger.info(f"Metadata extraction completed for: {filepath}")
        return result

    def _get_basic_info(self, filepath: Path) -> Dict[str, Any]:
        """Extract filesystem metadata"""
        try:
            self.logger.debug(f"Reading file stats for: {filepath}")
            stat = filepath.stat()
            mime_type, _ = mimetypes.guess_type(str(filepath))

            result = {
                "filename": filepath.name,
                "path": str(filepath.absolute()),
                "size_bytes": stat.st_size,
                "size_mb": round(stat.st_size / (1024 * 1024), 2),
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "extension": filepath.suffix.lower(),
                "mime_type": mime_type,
            }

            self.logger.debug(
                f"Basic info extracted: size={result['size_mb']}MB, mime_type={mime_type}"
            )
            return result

        except Exception as e:
            self.logger.error(f"Failed to read basic info for {filepath}: {e}")
            return {"error": f"Failed to read basic info: {e}"}

    def _extract_image_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract image metadata with lazy imports"""
        result = {"type": "image"}
        self.logger.debug("Starting image metadata extraction")

        # PIL metadata
        self.logger.debug("Attempting PIL metadata extraction")
        pil_data = self._get_pil_metadata(filepath)
        if "error" not in pil_data:
            self.logger.info("PIL metadata extraction successful")
            result.update(pil_data)
        else:
            self.logger.warning(
                f"PIL metadata extraction failed: {pil_data.get('error')}"
            )

        # Advanced metadata (optional)
        self.logger.debug("Attempting advanced metadata extraction")
        advanced_data = self._get_advanced_image_metadata(filepath)
        if "error" not in advanced_data:
            self.logger.info("Advanced metadata extraction successful")
            result.update(advanced_data)
        else:
            self.logger.debug(
                f"Advanced metadata extraction failed: {advanced_data.get('error')}"
            )

        return result

    def _get_pil_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract PIL-based image metadata"""
        try:
            self.logger.debug("Importing PIL dependencies")
            from PIL import Image
            from PIL.ExifTags import TAGS
            import pillow_heif

            pillow_heif.register_heif_opener()

            self.logger.debug(f"Opening image with PIL: {filepath}")
            with Image.open(filepath) as img:
                # Basic info
                image_info = {
                    "dimensions": img.size,
                    "width": img.size[0],
                    "height": img.size[1],
                    "mode": img.mode,
                    "format": img.format,
                    "has_transparency": img.mode in ("RGBA", "LA")
                    or "transparency" in img.info,
                }

                self.logger.debug(
                    f"Image dimensions: {img.size}, mode: {img.mode}, format: {img.format}"
                )

                # EXIF data
                exif_data = {}
                exif = img.getexif()  # Modern method, not deprecated _getexif
                if exif:
                    exif_data = {TAGS.get(k, f"tag_{k}"): v for k, v in exif.items()}
                    self.logger.debug(f"EXIF data extracted: {len(exif_data)} tags")
                else:
                    self.logger.debug("No EXIF data found")

                return {"image_info": image_info, "exif": exif_data}

        except ImportError as e:
            self.logger.error(f"PIL dependencies not available: {e}")
            return {"error": f"PIL not available: {e}"}
        except Exception as e:
            self.logger.error(f"PIL extraction failed for {filepath}: {e}")
            return {"error": f"PIL extraction failed: {e}"}

    def _get_advanced_image_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract advanced metadata using pyexiv2 if available"""
        try:
            self.logger.debug("Attempting pyexiv2 import")
            import pyexiv2

            self.logger.debug(f"Opening image with pyexiv2: {filepath}")
            with pyexiv2.Image(str(filepath)) as img:
                result = {
                    "advanced_metadata": {
                        "exif": img.read_exif(),
                        "iptc": img.read_iptc(),
                        "xmp": img.read_xmp(),
                    }
                }
                self.logger.info(
                    "Advanced metadata (EXIF/IPTC/XMP) extracted successfully"
                )
                return result

        except ImportError:
            self.logger.debug("pyexiv2 not available - using basic EXIF only")
            return {"info": "pyexiv2 not available - using basic EXIF only"}
        except Exception as e:
            self.logger.error(
                f"Advanced metadata extraction failed for {filepath}: {e}"
            )
            return {"error": f"Advanced metadata extraction failed: {e}"}

    def _extract_document_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Route document extraction based on file type"""
        result = {"type": "document"}
        ext = filepath.suffix.lower()

        extractors = {
            ".pdf": self._extract_pdf_metadata,
            ".docx": self._extract_docx_metadata,
            ".xlsx": self._extract_xlsx_metadata,
            ".pptx": self._extract_pptx_metadata,
        }

        if ext in extractors:
            self.logger.info(f"Using {ext} extractor")
            result.update(extractors[ext](filepath))
        else:
            self.logger.error(f"No extractor available for document type: {ext}")
            result["error"] = f"No extractor for {ext}"

        return result

    def _extract_pdf_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract PDF metadata - prefer PyMuPDF for comprehensive data"""
        self.logger.debug("Starting PDF metadata extraction")

        # Try PyMuPDF first (best overall)
        self.logger.debug("Attempting PyMuPDF extraction")
        pymupdf_result = self._try_pymupdf(filepath)
        if "error" not in pymupdf_result:
            self.logger.info("PDF metadata extracted successfully with PyMuPDF")
            return pymupdf_result

        # Fallback to pdfplumber for basic info
        self.logger.info("PyMuPDF failed, falling back to pdfplumber")
        result = self._try_pdfplumber(filepath)
        if "error" not in result:
            self.logger.info("PDF metadata extracted successfully with pdfplumber")
        else:
            self.logger.error("All PDF extraction methods failed")

        return result

    def _try_pymupdf(self, filepath: Path) -> Dict[str, Any]:
        """PyMuPDF extraction with proper resource management"""
        try:
            self.logger.debug("Importing PyMuPDF (fitz)")
            import fitz

            self.logger.debug(f"Opening PDF with PyMuPDF: {filepath}")
            doc = fitz.open(str(filepath))
            try:
                result = {
                    "pdf_info": {
                        "page_count": doc.page_count,
                        "is_encrypted": doc.needs_pass,
                        "metadata": doc.metadata,
                        "toc": doc.get_toc(),
                        "has_text": (
                            bool(doc[0].get_text()) if doc.page_count > 0 else False
                        ),
                    }
                }
                self.logger.debug(
                    f"PDF info: {doc.page_count} pages, encrypted: {doc.needs_pass}"
                )
                return result
            finally:
                doc.close()  # Ensure cleanup
                self.logger.debug("PyMuPDF document closed")

        except ImportError:
            self.logger.warning("PyMuPDF not available")
            return {"error": "PyMuPDF not available"}
        except Exception as e:
            self.logger.error(f"PyMuPDF extraction failed for {filepath}: {e}")
            return {"error": f"PyMuPDF extraction failed: {e}"}

    def _try_pdfplumber(self, filepath: Path) -> Dict[str, Any]:
        """Fallback PDF extraction"""
        try:
            self.logger.debug("Importing pdfplumber")
            import pdfplumber

            self.logger.debug(f"Opening PDF with pdfplumber: {filepath}")
            with pdfplumber.open(filepath) as pdf:
                result = {
                    "pdf_info": {
                        "page_count": len(pdf.pages),
                        "has_text": (
                            bool(pdf.pages[0].extract_text()) if pdf.pages else False
                        ),
                    }
                }
                self.logger.debug(f"PDF info: {len(pdf.pages)} pages")
                return result

        except ImportError:
            self.logger.error(
                "No PDF libraries available (neither PyMuPDF nor pdfplumber)"
            )
            return {"error": "No PDF libraries available"}
        except Exception as e:
            self.logger.error(f"pdfplumber extraction failed for {filepath}: {e}")
            return {"error": f"PDF extraction failed: {e}"}

    def _extract_docx_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract Word document metadata"""
        try:
            self.logger.debug("Importing python-docx")
            from docx import Document

            self.logger.debug(f"Opening DOCX document: {filepath}")
            doc = Document(str(filepath))
            props = doc.core_properties

            result = {
                "document_info": {
                    "title": props.title,
                    "author": props.author,
                    "created": props.created.isoformat() if props.created else None,
                    "modified": props.modified.isoformat() if props.modified else None,
                    "last_modified_by": props.last_modified_by,
                    "paragraph_count": len(doc.paragraphs),
                    "table_count": len(doc.tables),
                }
            }

            self.logger.info(
                f"DOCX metadata extracted: {len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables"
            )
            return result

        except ImportError:
            self.logger.error("python-docx library not available")
            return {"error": "python-docx not available"}
        except Exception as e:
            self.logger.error(f"DOCX extraction failed for {filepath}: {e}")
            return {"error": f"DOCX extraction failed: {e}"}

    def _extract_xlsx_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract Excel metadata"""
        try:
            self.logger.debug("Importing openpyxl")
            import openpyxl

            self.logger.debug(f"Opening XLSX workbook: {filepath}")
            wb = openpyxl.load_workbook(str(filepath), data_only=True)
            props = wb.properties

            result = {
                "spreadsheet_info": {
                    "title": props.title,
                    "creator": props.creator,
                    "created": props.created.isoformat() if props.created else None,
                    "modified": props.modified.isoformat() if props.modified else None,
                    "sheet_names": wb.sheetnames,
                    "sheet_count": len(wb.worksheets),
                }
            }

            self.logger.info(
                f"XLSX metadata extracted: {len(wb.worksheets)} sheets - {wb.sheetnames}"
            )
            return result

        except ImportError:
            self.logger.error("openpyxl library not available")
            return {"error": "openpyxl not available"}
        except Exception as e:
            self.logger.error(f"Excel extraction failed for {filepath}: {e}")
            return {"error": f"Excel extraction failed: {e}"}

    def _extract_pptx_metadata(self, filepath: Path) -> Dict[str, Any]:
        """Extract PowerPoint metadata"""
        try:
            self.logger.debug("Importing python-pptx")
            from pptx import Presentation

            self.logger.debug(f"Opening PPTX presentation: {filepath}")
            prs = Presentation(str(filepath))
            props = prs.core_properties

            result = {
                "presentation_info": {
                    "title": props.title,
                    "author": props.author,
                    "created": props.created.isoformat() if props.created else None,
                    "modified": props.modified.isoformat() if props.modified else None,
                    "slide_count": len(prs.slides),
                }
            }

            self.logger.info(f"PPTX metadata extracted: {len(prs.slides)} slides")
            return result

        except ImportError:
            self.logger.error("python-pptx library not available")
            return {"error": "python-pptx not available"}
        except Exception as e:
            self.logger.error(f"PowerPoint extraction failed for {filepath}: {e}")
            return {"error": f"PowerPoint extraction failed: {e}"}
